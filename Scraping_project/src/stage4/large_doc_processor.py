import logging
from datetime import datetime
from typing import Any

import httpx
from bs4 import BeautifulSoup
from tenacity import retry, stop_after_attempt, wait_exponential

from src.common.delta_lake import get_delta_manager

try:
    import sys
    from pathlib import Path

    sys.path.insert(0, str(Path(__file__).parent.parent.parent))
    from monitoring.metrics_exporter import (
        stage4_http_failures_total,
        stage4_http_requests_total,
    )

    PROMETHEUS_AVAILABLE = True
except ImportError:
    stage4_http_requests_total = None
    stage4_http_failures_total = None
    PROMETHEUS_AVAILABLE = False

logger = logging.getLogger(__name__)

class LargeDocProcessor:

    def __init__(self, model_name: str = "facebook/bart-large-cnn"):
        self.delta = get_delta_manager()
        self.model_name = model_name
        self.summarizer: Any = None

        self.CHUNK_SIZE = 5000
        self.OVERLAP = 500

        self.http_client = httpx.Client(
            headers={"User-Agent": "MyScraper/1.0 (Educational Research Bot)"},
            timeout=httpx.Timeout(30.0),
            follow_redirects=True,
        )

    def __del__(self):
        if hasattr(self, "http_client"):
            self.http_client.close()

    def _load_model(self):
        if self.summarizer is not None:
            return

        try:
            from transformers import pipeline

            logger.info(f"Loading heavyweight model: {self.model_name}")
            self.summarizer = pipeline(
                "summarization",
                model=self.model_name,
                device=-1,
            )
            logger.info("Model loaded successfully")
        except Exception as e:
            logger.error(f"Failed to load model: {e}")
            raise

    @retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
    def _fetch_content(self, url: str, is_pdf: bool = False) -> tuple[str, str]:
        try:
            if PROMETHEUS_AVAILABLE and stage4_http_requests_total:
                stage4_http_requests_total.inc()

            response = self.http_client.get(url)
            response.raise_for_status()

            content_type = response.headers.get("Content-Type", "").lower()

            if "application/pdf" in content_type or is_pdf:
                return self._extract_pdf_text(response.content), "pdf"

            elif (
                "application/vnd.openxmlformats-officedocument.wordprocessingml" in content_type
                or url.lower().endswith(".docx")
            ):
                return self._extract_docx_text(response.content), "docx"

            elif "application/vnd.openxmlformats-officedocument.presentationml" in content_type or url.lower().endswith(
                ".pptx"
            ):
                return self._extract_pptx_text(response.content), "pptx"

            elif "application/vnd.openxmlformats-officedocument.spreadsheetml" in content_type or url.lower().endswith(
                ".xlsx"
            ):
                return self._extract_xlsx_text(response.content), "xlsx"

            elif "application/msword" in content_type or url.lower().endswith(".doc"):
                return self._extract_doc_text(response.content), "doc"

            elif "text/html" in content_type:
                return self._extract_html_text(response.text), "html"

            elif "text/plain" in content_type:
                return response.text, "txt"

            else:
                logger.warning(f"Unsupported content type: {content_type} for {url}")
                return "", "unknown"

        except httpx.HTTPStatusError as e:
            if PROMETHEUS_AVAILABLE and stage4_http_failures_total:
                stage4_http_failures_total.labels(error_type="HTTPStatusError").inc()
            logger.error(f"HTTP error fetching {url}: {e.response.status_code}")
            raise
        except httpx.RequestError as e:
            if PROMETHEUS_AVAILABLE and stage4_http_failures_total:
                stage4_http_failures_total.labels(error_type="RequestError").inc()
            logger.error(f"Request error fetching {url}: {e}")
            raise
        except Exception as e:
            if PROMETHEUS_AVAILABLE and stage4_http_failures_total:
                stage4_http_failures_total.labels(error_type="UnknownError").inc()
            logger.error(f"Unexpected error fetching {url}: {e}")
            raise

    def _extract_html_text(self, html: str) -> str:
        try:
            soup = BeautifulSoup(html, "html.parser")

            for tag in soup(["script", "style", "nav", "header", "footer", "aside", "iframe"]):
                tag.decompose()

            text = soup.get_text(separator=" ", strip=True)
            text = " ".join(text.split())

            return text
        except Exception as e:
            logger.error(f"Failed to extract HTML text: {e}")
            return ""

    def _extract_pdf_text(self, pdf_content: bytes) -> str:
        try:
            from io import BytesIO

            from pypdf import PdfReader

            pdf_file = BytesIO(pdf_content)
            reader = PdfReader(pdf_file)

            text_parts = []
            for page in reader.pages:
                text_parts.append(page.extract_text())

            return " ".join(text_parts)

        except ImportError:
            logger.error("pypdf not installed - cannot extract PDF text")
            return ""
        except Exception as e:
            logger.error(f"Failed to extract PDF text: {e}")
            return ""

    def _extract_docx_text(self, docx_content: bytes) -> str:
        try:
            from io import BytesIO

            from docx import Document

            docx_file = BytesIO(docx_content)
            doc = Document(docx_file)

            text_parts = []
            for paragraph in doc.paragraphs:
                text_parts.append(paragraph.text)

            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_parts.append(cell.text)

            return "\n".join(text_parts)

        except ImportError:
            logger.error("python-docx not installed - cannot extract DOCX text")
            return ""
        except Exception as e:
            logger.error(f"Failed to extract DOCX text: {e}")
            return ""

    def _extract_pptx_text(self, pptx_content: bytes) -> str:
        try:
            from io import BytesIO

            from pptx import Presentation

            pptx_file = BytesIO(pptx_content)
            prs = Presentation(pptx_file)

            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)

            return "\n".join(text_parts)

        except ImportError:
            logger.error("python-pptx not installed - cannot extract PPTX text")
            return ""
        except Exception as e:
            logger.error(f"Failed to extract PPTX text: {e}")
            return ""

    def _extract_xlsx_text(self, xlsx_content: bytes) -> str:
        try:
            from io import BytesIO

            from openpyxl import load_workbook

            xlsx_file = BytesIO(xlsx_content)
            wb = load_workbook(xlsx_file, data_only=True)

            text_parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text_parts.append(str(cell.value))

            return "\n".join(text_parts)

        except ImportError:
            logger.error("openpyxl not installed - cannot extract XLSX text")
            return ""
        except Exception as e:
            logger.error(f"Failed to extract XLSX text: {e}")
            return ""

    def _extract_doc_text(self, doc_content: bytes) -> str:
        try:
            import tempfile

            import textract

            with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
                tmp.write(doc_content)
                tmp_path = tmp.name

            text = textract.process(tmp_path).decode("utf-8")

            import os

            os.unlink(tmp_path)

            return text

        except ImportError:
            logger.error("textract not installed - cannot extract .doc text. Install: apt-get install antiword")
            return ""
        except Exception as e:
            logger.error(f"Failed to extract .doc text: {e}")
            return ""

    def process_large_document(self, url: str, text: str) -> str:
        try:
            chunks = self._split_into_chunks(text)
            logger.info(f"Split {url[:80]} into {len(chunks)} chunks")

            chunk_summaries = []
            for i, chunk in enumerate(chunks):
                try:
                    summary = self._summarize_chunk(chunk)
                    if summary:
                        chunk_summaries.append(summary)
                except Exception as e:
                    logger.warning(f"Failed to summarize chunk {i}: {e}")

            if not chunk_summaries:
                return text[:500] + "..." if len(text) > 500 else text

            combined_summary = " ".join(chunk_summaries)

            if len(combined_summary) > 1000:
                refined_summary = self._summarize_chunk(combined_summary[:5000])
                if refined_summary:
                    combined_summary = refined_summary

            return combined_summary

        except Exception as e:
            logger.error(f"Failed to process large document: {e}")
            return text[:500] + "..." if len(text) > 500 else text

    def process_queue(self):
        try:
            all_docs = self.delta.read("stage4_large_docs")
            pending_docs = [d for d in all_docs if d.get("status") == "pending"]

            if not pending_docs:
                logger.info("No pending large documents to process")
                return

            logger.info(f"Processing {len(pending_docs)} large documents")

            self._load_model()

            summaries = []
            processed_count = 0

            for doc in pending_docs:
                try:
                    summary = self._process_document(doc)
                    if summary:
                        summaries.append(summary)
                        processed_count += 1
                except Exception as e:
                    logger.error(f"Failed to process doc {doc.get('url')}: {e}")

            if summaries:
                self.delta.write("stage4_summaries", summaries, mode="append", async_write=False)
                logger.info(f"Saved {len(summaries)} summaries")

            self._update_queue_status(all_docs, pending_docs)

            logger.info(f"Completed processing {processed_count} large documents")

        except Exception as e:
            logger.error(f"Queue processing failed: {e}", exc_info=True)

    def _process_document(self, doc: dict[str, Any]) -> dict[str, Any] | None:
        url_value = doc.get("url")
        if not isinstance(url_value, str):
            logger.warning("Document missing URL; skipping entry")
            return None

        url = url_value
        is_pdf = bool(doc.get("is_pdf", False))

        logger.info(f"Processing large doc: {url[:80]}")

        try:
            text, content_type = self._fetch_content(url, is_pdf)

            if not text:
                logger.warning(f"No text extracted from {url}")
                return None

            word_count = len(text.split())
            logger.info(f"Fetched {word_count} words from {url[:80]} (type: {content_type})")

        except Exception as e:
            logger.error(f"Failed to fetch content from {url}: {e}")
            return None

        chunks = self._split_into_chunks(text)
        logger.info(f"Split into {len(chunks)} chunks")

        chunk_summaries = []
        for i, chunk in enumerate(chunks):
            try:
                summary = self._summarize_chunk(chunk)
                if summary:
                    chunk_summaries.append(summary)
                    logger.debug(f"Summarized chunk {i + 1}/{len(chunks)}")
            except Exception as e:
                logger.warning(f"Failed to summarize chunk {i}: {e}")

        if not chunk_summaries:
            logger.warning(f"No summaries generated for {url}")
            return None

        combined_summary = " ".join(chunk_summaries)

        if len(combined_summary) > 1000:
            refined_summary = self._summarize_chunk(combined_summary[:5000])
            if refined_summary:
                combined_summary = refined_summary

        return {
            "url": url,
            "summary": combined_summary,
            "original_word_count": word_count,
            "chunk_count": len(chunks),
            "processed_at": datetime.now().isoformat(),
            "model_used": self.model_name,
        }

    def _split_into_chunks(self, text: str) -> list[str]:
        if len(text) <= self.CHUNK_SIZE:
            return [text]

        chunks = []
        start = 0

        while start < len(text):
            end = start + self.CHUNK_SIZE
            chunk = text[start:end]

            if end < len(text):
                last_period = chunk.rfind(".")
                if last_period > self.CHUNK_SIZE // 2:
                    end = start + last_period + 1
                    chunk = text[start:end]

            chunks.append(chunk.strip())
            start = end - self.OVERLAP

        return chunks

    def _summarize_chunk(self, text: str) -> str | None:
        if not text or len(text) < 100:
            return None

        try:
            max_input = 1024
            if len(text) > max_input:
                text = text[:max_input]

            result = self.summarizer(text, max_length=150, min_length=30, do_sample=False)

            return result[0]["summary_text"]

        except Exception as e:
            logger.error(f"Chunk summarization failed: {e}")
            sentences = text.split(".")[:3]
            return ". ".join(sentences) + "."

    def _update_queue_status(
        self,
        all_docs: list[dict[str, Any]],
        processed_docs: list[dict[str, Any]],
    ) -> None:
        """Update queue with completed status."""
        processed_urls = {d.get("url") for d in processed_docs}

        updated_queue = []
        for doc in all_docs:
            if doc.get("url") in processed_urls:
                doc["status"] = "completed"
                doc["completed_at"] = datetime.now().isoformat()
            updated_queue.append(doc)

        self.delta.write("stage4_large_docs", updated_queue, mode="overwrite", async_write=False)
        logger.info(f"Updated queue status for {len(processed_urls)} documents")

def process_large_documents():
    processor = LargeDocProcessor()
    processor.process_queue()

if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(name)s: %(message)s")
    process_large_documents()
